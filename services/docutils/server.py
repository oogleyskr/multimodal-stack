"""
Document Utilities microservice for parsing various document formats.

CPU-only service — no GPU required. Uses pymupdf, python-docx, openpyxl,
python-pptx, and beautifulsoup4 to extract text from documents.

Port: 8106
"""

import io
import os
import time
import logging
import tempfile

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("docutils")

app = FastAPI(title="Local Document Utils Service", version="1.0.0")


@app.get("/health")
async def health():
    """Health check — always ready (CPU-only, no model to load)."""
    return {"status": "ok", "service": "docutils"}


def parse_pdf(file_path: str) -> dict:
    """Extract text and metadata from a PDF file using pymupdf."""
    import fitz  # pymupdf

    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        pages.append({"page": i + 1, "text": text.strip()})

    metadata = doc.metadata or {}
    doc.close()

    return {
        "format": "pdf",
        "pages": len(pages),
        "metadata": {k: v for k, v in metadata.items() if v},
        "content": pages,
        "full_text": "\n\n".join(p["text"] for p in pages if p["text"]),
    }


def parse_docx(file_path: str) -> dict:
    """Extract text from a Word document (.docx)."""
    from docx import Document

    doc = Document(file_path)

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        tables.append(rows)

    return {
        "format": "docx",
        "paragraphs": len(paragraphs),
        "tables": len(tables),
        "content": paragraphs,
        "table_data": tables,
        "full_text": "\n\n".join(paragraphs),
    }


def parse_xlsx(file_path: str) -> dict:
    """Extract data from an Excel spreadsheet (.xlsx)."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    sheets = {}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            # Convert each cell to string, skip completely empty rows
            str_row = [str(cell) if cell is not None else "" for cell in row]
            if any(c.strip() for c in str_row):
                rows.append(str_row)
        sheets[sheet_name] = rows

    wb.close()

    return {
        "format": "xlsx",
        "sheets": list(sheets.keys()),
        "sheet_count": len(sheets),
        "data": sheets,
        "full_text": "\n".join(
            f"[{name}]\n" + "\n".join("\t".join(r) for r in rows)
            for name, rows in sheets.items()
        ),
    }


def parse_pptx(file_path: str) -> dict:
    """Extract text from a PowerPoint presentation (.pptx)."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        slides.append({"slide": i + 1, "text": texts})

    return {
        "format": "pptx",
        "slides": len(slides),
        "content": slides,
        "full_text": "\n\n".join(
            f"--- Slide {s['slide']} ---\n" + "\n".join(s["text"])
            for s in slides if s["text"]
        ),
    }


def parse_html(file_path: str) -> dict:
    """Extract text from an HTML file."""
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()

    title = soup.title.string if soup.title else ""
    text = soup.get_text(separator="\n", strip=True)

    return {
        "format": "html",
        "title": title,
        "full_text": text,
    }


def parse_text(file_path: str) -> dict:
    """Read a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    return {
        "format": "text",
        "full_text": text,
        "lines": text.count("\n") + 1,
        "characters": len(text),
    }


# Map file extensions to parser functions
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".txt": parse_text,
    ".md": parse_text,
    ".csv": parse_text,
    ".json": parse_text,
    ".xml": parse_text,
    ".yaml": parse_text,
    ".yml": parse_text,
    ".log": parse_text,
}


@app.post("/parse")
async def parse_document(
    file: UploadFile = File(...),
    pages: str = Form(default="", description="Page range for PDFs, e.g. '1-5' (optional)"),
):
    """
    Parse a document and extract its text content.

    Supported formats: PDF, DOCX, XLSX, PPTX, HTML, TXT, MD, CSV, JSON, XML, YAML.

    Args:
        file: The document file to parse.
        pages: Optional page range for PDFs (e.g. "1-5", "3,7,10-12").

    Returns:
        JSON with extracted text content and metadata.
    """
    filename = file.filename or "document"
    ext = os.path.splitext(filename)[1].lower()

    if ext not in PARSERS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported format: {ext}. Supported: {', '.join(sorted(PARSERS.keys()))}"
        )

    # Write to temp file for parsing
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = tmp.name

    try:
        start = time.time()
        result = PARSERS[ext](tmp_path)
        elapsed = time.time() - start

        result["filename"] = filename
        result["file_size"] = len(content)
        result["processing_time"] = round(elapsed, 3)

        logger.info(f"Parsed {filename} ({len(content)} bytes) in {elapsed:.1f}s")

        return JSONResponse(result)
    except Exception as e:
        logger.error(f"Failed to parse {filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse {filename}: {str(e)}")
    finally:
        os.unlink(tmp_path)


@app.get("/formats")
async def supported_formats():
    """List all supported document formats."""
    return {"formats": sorted(PARSERS.keys())}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8106)
